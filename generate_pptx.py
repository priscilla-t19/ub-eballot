from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

UB_BLUE   = RGBColor(0, 48, 135)
UB_GOLD   = RGBColor(255, 215, 0)
WHITE     = RGBColor(255, 255, 255)
LIGHT     = RGBColor(240, 244, 255)
DARK_GRAY = RGBColor(50, 50, 50)
MID_GRAY  = RGBColor(110, 110, 110)
SUCCESS   = RGBColor(26, 122, 74)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def add_rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=DARK_GRAY,
             align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.3, fill=UB_BLUE)
    add_rect(slide, 0, 1.3, 13.33, 0.07, fill=UB_GOLD)
    add_text(slide, title, 0.4, 0.1, 12.5, 0.72, size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.82, 12.5, 0.42, size=13,
                 color=RGBColor(180, 200, 255), italic=True)


def footer(slide, text='UB eBallot  —  Beta Presentation  |  University of Botswana  |  2026'):
    add_rect(slide, 0, 7.15, 13.33, 0.35, fill=UB_BLUE)
    add_text(slide, text, 0.3, 7.18, 12.7, 0.28,
             size=10, color=RGBColor(180, 200, 255), align=PP_ALIGN.CENTER)


# ── SLIDE 1  Title ───────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, fill=UB_BLUE)
add_rect(s, 0, 5.8, 13.33, 0.14, fill=UB_GOLD)
add_text(s, 'UB', 0.6, 1.1, 3, 1.5, size=100, bold=True, color=UB_GOLD)
add_text(s, 'eBallot', 2.85, 1.5, 8, 1.1, size=70, bold=True, color=WHITE)
add_text(s, 'Secure Online Voting for UB SRC Elections',
         0.6, 3.25, 12, 0.6, size=22, color=RGBColor(180, 210, 255))
add_text(s, 'Beta Version Presentation  \u2022  University of Botswana  \u2022  2026',
         0.6, 3.95, 12, 0.5, size=14, color=RGBColor(140, 170, 220), italic=True)
add_text(s, 'Built with Flask \u00b7 PostgreSQL \u00b7 AES-128 Encryption \u00b7 HMAC-SHA256',
         0.6, 4.9, 12, 0.45, size=12, color=RGBColor(120, 150, 200), italic=True)


# ── SLIDE 2  The Problem ─────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, fill=LIGHT)
header_bar(s, 'The Problem',
           'Why UB SRC elections needed a digital solution')
footer(s)

problems = [
    ('Paper ballots are slow',
     'Counting is manual, error-prone, and takes hours after polling closes.'),
    ('No way to verify your vote',
     'Once a ballot is in the box there is no receipt and no way to confirm it was counted.'),
    ('Voter identity is hard to enforce',
     'Preventing impersonation or double-voting relies entirely on manual checks.'),
    ('Low turnout',
     'Students must be physically present on campus — a barrier for many eligible voters.'),
    ('Results are opaque',
     'The counting process is not independently verifiable by the student body.'),
]

y = 1.58
for title, desc in problems:
    add_rect(s, 0.4, y, 12.5, 0.88, fill=WHITE, line=RGBColor(200, 210, 230))
    add_text(s, '\u2717  ' + title, 0.65, y + 0.05, 5.2, 0.4,
             size=13, bold=True, color=RGBColor(180, 30, 30))
    add_text(s, desc, 5.6, y + 0.05, 7.1, 0.75, size=12, color=MID_GRAY)
    y += 1.0


# ── SLIDE 3  The Solution ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, fill=LIGHT)
header_bar(s, 'The Solution  \u2014  UB eBallot',
           'A secure, anonymous, and verifiable online voting platform')
footer(s)

cols = [
    (UB_BLUE,
     'Secure',
     'End-to-end AES-128 encryption. Student numbers are never stored alongside votes. Every admin action is logged in an audit trail.'),
    (SUCCESS,
     'Anonymous',
     'Votes are separated from identity using a one-way HMAC-SHA256 token. Not even administrators can link a ballot to a specific voter.'),
    (RGBColor(180, 100, 0),
     'Verifiable',
     'Every voter receives a cryptographic receipt code immediately after voting. Anyone can independently confirm their vote was counted.'),
    (RGBColor(100, 30, 160),
     'Accessible',
     'Fully web-based and mobile-friendly. Students vote from any device and any location during the open election window.'),
]

x = 0.35
for color, title, desc in cols:
    add_rect(s, x, 1.6, 2.9, 5.45, fill=WHITE, line=color)
    add_rect(s, x, 1.6, 2.9, 0.45, fill=color)
    add_text(s, title, x + 0.12, 1.65, 2.66, 0.38,
             size=16, bold=True, color=WHITE)
    add_text(s, desc, x + 0.15, 2.2, 2.6, 4.6, size=12, color=DARK_GRAY)
    x += 3.18


# ── SLIDE 4  Voter Flow ──────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, fill=LIGHT)
header_bar(s, 'How It Works  \u2014  Student Voting Flow', '')
footer(s)

steps = [
    (UB_BLUE,               '1', 'Visit the site',
     'Open the UB eBallot URL on any browser or smartphone.'),
    (UB_GOLD,               '2', 'Enter student number',
     'Type your UB student number. The system verifies it against the registry.'),
    (SUCCESS,               '3', 'Verify via email OTP',
     'A 6-digit code is sent to your UB email. Enter it to prove your identity.'),
    (RGBColor(0, 130, 180), '4', 'Cast your ballot',
     'Select your preferred candidates for each position and submit.'),
    (RGBColor(140, 30, 140),'5', 'Save your receipt',
     'A unique receipt code is displayed. Keep it to verify your vote later.'),
]

x = 0.28
for color, num, title, desc in steps:
    add_rect(s, x, 1.6, 2.45, 5.0, fill=WHITE, line=color)
    circ = s.shapes.add_shape(9, Inches(x + 0.9), Inches(1.78),
                              Inches(0.65), Inches(0.65))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    txt_col = UB_BLUE if color == UB_GOLD else WHITE
    add_text(s, num, x + 0.9, 1.8, 0.65, 0.55,
             size=20, bold=True, color=txt_col, align=PP_ALIGN.CENTER)
    add_text(s, title, x + 0.1, 2.6, 2.25, 0.48,
             size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s, desc,  x + 0.15, 3.18, 2.15, 3.2,
             size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)
    x += 2.6


# ── SLIDE 5  Admin Panel ─────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, fill=LIGHT)
header_bar(s, 'Admin Panel Features',
           'Full election management from a secure, role-protected dashboard')
footer(s)

features = [
    ('Election Management', [
        'Create elections with open / close times',
        'Activate or deactivate elections live',
        'Add positions and candidates',
        'Publish results when counting is done',
    ]),
    ('Student Registry', [
        'Import all students from a single CSV file',
        'Registry persists across all future elections',
        'Add or deactivate individual entries manually',
        'No student ID is exposed in any ballot view',
    ]),
    ('Database View', [
        'Live ballot count per election',
        'Voter turnout statistics at a glance',
        'Full audit log of every admin action',
        'Privacy confirmation panel built in',
    ]),
    ('Security Controls', [
        'Argon2id password hashing for the admin account',
        'Rate limiting on login and voting routes',
        'Session expires after 30 minutes of inactivity',
        'Dropdown nav reduces accidental exposure',
    ]),
]

x = 0.35
for title, items in features:
    add_rect(s, x, 1.6, 2.9, 5.55, fill=WHITE, line=UB_BLUE)
    add_rect(s, x, 1.6, 2.9, 0.44, fill=UB_BLUE)
    add_text(s, title, x + 0.12, 1.65, 2.66, 0.38,
             size=13, bold=True, color=WHITE)
    y2 = 2.18
    for item in items:
        add_text(s, '\u25b8  ' + item, x + 0.18, y2, 2.52, 0.55,
                 size=11, color=DARK_GRAY)
        y2 += 0.74
    x += 3.18


# ── SLIDE 6  Security Architecture ──────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, fill=LIGHT)
header_bar(s, 'Security Architecture',
           'How anonymity and ballot integrity are technically guaranteed')
footer(s)

add_rect(s, 0.35, 1.55, 6.0, 5.6, fill=WHITE, line=UB_BLUE)
add_rect(s, 0.35, 1.55, 6.0, 0.44, fill=UB_BLUE)
add_text(s, 'Identity \u2194 Vote Separation', 0.5, 1.6, 5.7, 0.38,
         size=14, bold=True, color=WHITE)
left_points = [
    'Student number enters the Auth module only',
    'Auth computes HMAC-SHA256(pepper, student_no:election_id)',
    'Only the anonymous token crosses into the voting module',
    'Student number is deleted from session immediately',
    'Ballot stored with no link to any student identity',
    'Even if the database is stolen, votes cannot be traced back',
]
y2 = 2.1
for p in left_points:
    add_text(s, '\u2713  ' + p, 0.55, y2, 5.6, 0.5, size=11, color=DARK_GRAY)
    y2 += 0.72

add_rect(s, 6.98, 1.55, 6.0, 5.6, fill=WHITE, line=SUCCESS)
add_rect(s, 6.98, 1.55, 6.0, 0.44, fill=SUCCESS)
add_text(s, 'Ballot Integrity', 7.12, 1.6, 5.7, 0.38,
         size=14, bold=True, color=WHITE)
right_points = [
    'Ballots encrypted with AES-128 (Fernet) before storage',
    'Each ballot has a SHA-256 chain hash linking to previous',
    'Any tampering with stored ballots immediately breaks the chain',
    'Admin can verify chain integrity live from the dashboard',
    'Receipt code is SHA-256 hashed — only the voter knows the raw code',
    'Verification proves inclusion without revealing who you voted for',
]
y2 = 2.1
for p in right_points:
    add_text(s, '\u2713  ' + p, 7.15, y2, 5.6, 0.5, size=11, color=DARK_GRAY)
    y2 += 0.72


# ── SLIDE 7  Technology Stack ────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, fill=LIGHT)
header_bar(s, 'Technology Stack', '')
footer(s)

tech = [
    ('Backend',      UB_BLUE,
     ['Python 3 / Flask', 'SQLAlchemy ORM', 'Flask-Mail (OTP delivery)', 'Gunicorn (production WSGI)']),
    ('Database',     SUCCESS,
     ['PostgreSQL on Railway', 'SQLite for local dev', 'Persistent student registry', 'Tamper-evident ballot chain']),
    ('Cryptography', RGBColor(140, 30, 140),
     ['AES-128 Fernet (ballots)', 'HMAC-SHA256 (voter tokens)', 'Argon2id (admin passwords)', 'SHA-256 (OTPs + receipts)']),
    ('Frontend',     RGBColor(180, 100, 0),
     ['Responsive HTML / CSS', 'Vanilla JavaScript', 'Progressive Web App (PWA)', 'Mobile-first design']),
    ('Deployment',   RGBColor(0, 130, 180),
     ['Railway (cloud hosting)', 'GitHub (version control)', 'Auto-deploy on git push', 'Environment-based config']),
]

x = 0.2
for title, color, items in tech:
    add_rect(s, x, 1.55, 2.45, 5.6, fill=WHITE, line=color)
    add_rect(s, x, 1.55, 2.45, 0.44, fill=color)
    add_text(s, title, x + 0.12, 1.6, 2.22, 0.38,
             size=13, bold=True, color=WHITE)
    y2 = 2.12
    for item in items:
        add_text(s, '\u25b8  ' + item, x + 0.18, y2, 2.1, 0.55,
                 size=11, color=DARK_GRAY)
        y2 += 0.78
    x += 2.62


# ── SLIDE 8  Beta Status & Roadmap ──────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, fill=LIGHT)
header_bar(s, 'Beta Status & Roadmap', '')
footer(s)

add_rect(s, 0.35, 1.55, 5.9, 5.6, fill=WHITE, line=SUCCESS)
add_rect(s, 0.35, 1.55, 5.9, 0.44, fill=SUCCESS)
add_text(s, '\u2713  Completed in Beta', 0.5, 1.6, 5.6, 0.38,
         size=14, bold=True, color=WHITE)
done = [
    'Student authentication via student number + OTP',
    'Anonymous ballot submission with AES-128 encryption',
    'Tamper-evident SHA-256 ballot chain',
    'Vote receipt and independent verification page',
    'Full admin dashboard (elections, positions, candidates)',
    'Persistent student registry with CSV import',
    'Rate limiting and brute-force protection',
    'Mobile-responsive Progressive Web App interface',
    'Cloud deployment on Railway with PostgreSQL',
]
y2 = 2.12
for item in done:
    add_text(s, '\u2713  ' + item, 0.55, y2, 5.5, 0.5, size=11, color=DARK_GRAY)
    y2 += 0.57

add_rect(s, 7.08, 1.55, 5.9, 5.6, fill=WHITE, line=UB_GOLD)
add_rect(s, 7.08, 1.55, 5.9, 0.44, fill=RGBColor(200, 160, 0))
add_text(s, '\u2192  Production Roadmap', 7.22, 1.6, 5.6, 0.38,
         size=14, bold=True, color=WHITE)
todo = [
    'Direct integration with UB student database (API)',
    'SMS OTP fallback for students without email access',
    'Two-factor authentication for admin login',
    'Automated daily database backups with alerts',
    'Load testing for the full UB student body (~18 000)',
    'Independent security audit / penetration test',
    'Printable results report for official SRC records',
    'Candidate photo upload and manifesto viewer',
]
y2 = 2.12
for item in todo:
    add_text(s, '\u2192  ' + item, 7.25, y2, 5.5, 0.5, size=11, color=DARK_GRAY)
    y2 += 0.64


# ── SLIDE 9  Q&A ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, fill=LIGHT)
header_bar(s, 'Anticipated Questions', '')
footer(s)

qas = [
    ('How do you ensure one vote per student?',
     'An HMAC-SHA256 token is derived from the student number + election ID. It is stored with a unique DB constraint — a second attempt is rejected even under simultaneous submissions (race-condition safe).'),
    ('Can the admin see who voted for whom?',
     'No. The student number never enters the voting module. Ballots are encrypted and linked only to an anonymous token. The admin sees counts only, not individual choices.'),
    ('What if the server goes down mid-vote?',
     'The ballot and voter token are committed in a single atomic database transaction. If either part fails, both are rolled back — no partial or lost votes.'),
    ('How does receipt verification work?',
     'After voting a random receipt code is shown to the student. Its SHA-256 hash is stored with the ballot. The verify page hashes the entered code and confirms it matches a record.'),
    ('Is the system vulnerable to hacking?',
     'Rate limiting, CSRF protection, and input sanitisation are applied on all routes. Passwords use Argon2id. Sessions expire after 30 min. The codebase follows OWASP Top-10 guidelines.'),
    ('Why not use ElectionBuddy or a commercial platform?',
     'Commercial tools have subscription costs, are not tailored to UB student numbers, and store data on foreign servers. UB eBallot is purpose-built, free, and data stays under UB control.'),
]

y = 1.55
for q, a in qas:
    add_rect(s, 0.35, y, 12.6, 0.85, fill=WHITE, line=UB_BLUE)
    add_text(s, 'Q:  ' + q, 0.55, y + 0.04, 12.2, 0.36,
             size=11.5, bold=True, color=UB_BLUE)
    add_text(s, 'A:  ' + a, 0.55, y + 0.42, 12.2, 0.38,
             size=10.5, color=DARK_GRAY)
    y += 0.94


# ── SLIDE 10  Thank You ──────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, fill=UB_BLUE)
add_rect(s, 0, 3.55, 13.33, 0.12, fill=UB_GOLD)
add_text(s, 'Thank You', 0, 1.15, 13.33, 1.5,
         size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 'Questions & Live Demo', 0, 2.85, 13.33, 0.7,
         size=24, color=UB_GOLD, align=PP_ALIGN.CENTER)
add_text(s, 'UB eBallot  \u2014  Secure  \u00b7  Anonymous  \u00b7  Verifiable',
         0, 4.15, 13.33, 0.6, size=16,
         color=RGBColor(180, 210, 255), align=PP_ALIGN.CENTER, italic=True)
add_text(s, 'Deployed on Railway  \u2022  Source: github.com/priscilla-t19/ub-eballot',
         0, 4.95, 13.33, 0.5, size=12,
         color=RGBColor(140, 170, 220), align=PP_ALIGN.CENTER)


prs.save('UB_eBallot_Beta_Presentation.pptx')
print('Saved: UB_eBallot_Beta_Presentation.pptx')
